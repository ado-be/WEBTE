# pdf_to_pptx.py
import sys
import os
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches

def pdf_to_pptx(pdf_path, output_path="output.pptx"):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # prázdny slide

    output_folder = "output_images"
    os.makedirs(output_folder, exist_ok=True)

    doc = fitz.open(pdf_path)
    for i, page in enumerate(doc):
        pix = page.get_pixmap(dpi=150)
        img_path = os.path.join(output_folder, f"slide_{i+1}.png")
        pix.save(img_path)

        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(10), Inches(7.5))

        os.remove(img_path)

    prs.save(output_path)
    print(f"PowerPoint uložený ako: {output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 2 or len(sys.argv) > 3:
        print("Použitie: python pdf_to_pptx.py vstup.pdf [vystup.pptx]")
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) == 3 else "output.pptx"
    pdf_to_pptx(input_file, output_file)
